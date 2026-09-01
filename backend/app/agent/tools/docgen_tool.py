"""Document Deliverable Generator — Phase 7.5 (Plan.md §7.5).

Generates real file deliverables (DOCX, XLSX, PPTX) from structured agent
output.  No model is involved here — this is pure document-assembly logic.

Supported output types
----------------------
``generate_approval_note_docx``
    The flagship deliverable: a formatted Word document following the 9-section
    structure from the document-deliverable-generator SKILL.md:
        1. Title
        2. Background
        3. Inspection Summary
        4. Key Findings
        5. SOP / Manual Evidence
        6. Risk & Impact
        7. Recommendation
        8. Human Approval Section  (signature block)
        9. Evidence Table          (citations from the RAG pipeline)

``generate_calculation_sheet_xlsx``
    An Excel workbook for engineering calculations (measurements vs. limits)
    that the agent surfaces alongside an approval note.

``generate_summary_deck_pptx``
    A minimal slide deck: title slide + one slide per key finding +
    a conclusion/recommendation slide.

All outputs are written to ``data/deliverables/<run_id>/`` and the full
path is returned to the caller so it can be served / linked in the UI.

SKILL constraints enforced
--------------------------
- Do not invent measurements.  All numeric fields must come from caller-
  supplied ``findings``; if a field is missing it is rendered as
  ``[PENDING — source data required]``.
- Do not invent dates.  Dates are either passed in or stamped with today's
  UTC date prefixed ``[AUTO]``.
- Include evidence references where possible.  Every finding row in the
  DOCX evidence table cites source_file, page_number, and exact_quote from
  a RetrievalResult.

Sovereignty: all libraries (python-docx, openpyxl, python-pptx) run
fully locally.  No external calls.
"""

from __future__ import annotations

import logging
import os
from dataclasses import dataclass, field
from datetime import date, datetime, timezone
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Output directory
# ---------------------------------------------------------------------------

# Resolved relative to the server's working directory (backend/).
_DELIVERABLES_ROOT = Path("data/deliverables")

PLACEHOLDER = "[PENDING — source data required]"


# ---------------------------------------------------------------------------
# Input data classes
# ---------------------------------------------------------------------------

@dataclass
class Finding:
    """One flagged observation to include in the approval note."""
    label: str                     # e.g. "F-1 (ALERT)"
    description: str               # full finding text
    measured_value: Optional[str] = None   # e.g. "6.7 mm/s RMS"
    limit_value: Optional[str] = None      # e.g. "6.0 mm/s (SOP-17 §3.2)"
    status: str = "ALERT"          # NORMAL | ALERT | FLAG | CAUTION | SHUTDOWN


@dataclass
class EvidenceCitation:
    """One RAG-retrieved chunk used as evidence. Maps to RetrievalResult."""
    source_file: str
    page_number: int
    exact_quote: str
    confidence: float
    doc_type: str = "other"
    section_title: Optional[str] = None


@dataclass
class ApprovalNoteInput:
    """All structured data the planner must supply to generate an approval note."""
    # --- Header ---
    subject: str                          # e.g. "Continued Operation — Pump P-204"
    equipment_id: str                     # e.g. "P-204"
    inspection_report_id: str            # e.g. "IR-892"
    inspection_date: str                 # e.g. "2026-08-15"
    inspector_name: str                  # e.g. "R. K. Sharma"
    # --- Content ---
    background: str                       # 1–3 sentence context paragraph
    findings: list[Finding] = field(default_factory=list)
    applicable_sop: str = PLACEHOLDER     # e.g. "SOP-17, Rev 4"
    risk_assessment: str = PLACEHOLDER    # paragraph(s)
    conditions: list[str] = field(default_factory=list)  # bullet list for §4
    recommendation: str = PLACEHOLDER
    # --- Evidence ---
    evidence: list[EvidenceCitation] = field(default_factory=list)
    # --- Optional overrides ---
    run_id: Optional[int] = None
    raised_date: Optional[str] = None    # ISO date; auto-stamped if None
    area_engineer_name: str = PLACEHOLDER
    note_id: Optional[str] = None        # e.g. "AN-AUTO-001"


@dataclass
class CalculationSheetInput:
    """Data for a measurement-vs-limit Excel sheet."""
    title: str
    equipment_id: str
    measurements: list[dict]   # each: {parameter, reading, unit, limit, status}
    run_id: Optional[int] = None


@dataclass
class SummaryDeckInput:
    """Data for a brief PPTX summary deck."""
    title: str
    equipment_id: str
    inspection_report_id: str
    findings: list[Finding] = field(default_factory=list)
    recommendation: str = PLACEHOLDER
    run_id: Optional[int] = None


# ---------------------------------------------------------------------------
# Output path helper
# ---------------------------------------------------------------------------

def _output_path(run_id: Optional[int], filename: str) -> Path:
    folder = _DELIVERABLES_ROOT / (str(run_id) if run_id is not None else "untracked")
    folder.mkdir(parents=True, exist_ok=True)
    return folder / filename


def _today() -> str:
    return date.today().isoformat()


def _stamp(provided: Optional[str]) -> str:
    return provided if provided else f"[AUTO] {_today()}"


# ---------------------------------------------------------------------------
# DOCX — Approval Note
# ---------------------------------------------------------------------------

def generate_approval_note_docx(data: ApprovalNoteInput) -> Path:
    """Generate a formatted DOCX approval note from *data*.

    Returns the absolute path to the created file.
    """
    from docx import Document
    from docx.shared import Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement

    doc = Document()

    # ---- Page margins (narrow) ----
    for section in doc.sections:
        section.top_margin    = _cm(2)
        section.bottom_margin = _cm(2)
        section.left_margin   = _cm(2.5)
        section.right_margin  = _cm(2.5)

    # ---- Title block ----
    note_id = data.note_id or f"AN-RUN-{data.run_id or 'DRAFT'}"
    stamp_date = _stamp(data.raised_date)

    title_para = doc.add_heading(f"APPROVAL NOTE — {note_id}", level=1)
    title_para.runs[0].font.color.rgb = RGBColor(0x1a, 0x37, 0x6e)

    _add_label_value(doc, "Subject",    data.subject)
    _add_label_value(doc, "Equipment",  data.equipment_id)
    _add_label_value(doc, "Date",       stamp_date)
    _add_label_value(doc, "Raised by",  data.inspector_name)
    _add_label_value(doc, "Reviewed by", data.area_engineer_name)

    doc.add_paragraph()  # spacer

    # ---- §1 Reference Documents ----
    doc.add_heading("1.  Reference Documents", level=2)
    _add_label_value(doc, "Inspection Report", data.inspection_report_id)
    _add_label_value(doc, "Inspection Date",   data.inspection_date)
    _add_label_value(doc, "Applicable SOP",    data.applicable_sop)

    # ---- §2 Background ----
    doc.add_heading("2.  Background", level=2)
    doc.add_paragraph(data.background or PLACEHOLDER)

    # ---- §3 Inspection Summary / Key Findings ----
    doc.add_heading("3.  Inspection Summary", level=2)
    if data.findings:
        table = doc.add_table(rows=1, cols=4)
        table.style = "Table Grid"
        hdr = table.rows[0].cells
        for cell, txt in zip(hdr, ["Finding", "Description", "Measured", "Limit / Status"]):
            cell.text = txt
            _bold_cell(cell)
        for f in data.findings:
            row = table.add_row().cells
            row[0].text = f.label
            row[1].text = f.description
            row[2].text = f.measured_value or PLACEHOLDER
            row[3].text = f"{f.limit_value or PLACEHOLDER}  [{f.status}]"
    else:
        doc.add_paragraph(PLACEHOLDER)

    # ---- §4 SOP / Manual Evidence ----
    doc.add_heading("4.  SOP / Manual Evidence", level=2)
    doc.add_paragraph(
        f"Applicable standard: {data.applicable_sop}. "
        "Key clauses cited in the findings above. "
        "Full evidence quotes are provided in the Evidence Table (§9)."
    )

    # ---- §5 Risk & Impact ----
    doc.add_heading("5.  Risk & Impact", level=2)
    doc.add_paragraph(data.risk_assessment)

    # ---- §6 Conditions for Continued Operation ----
    doc.add_heading("6.  Conditions for Continued Operation", level=2)
    if data.conditions:
        for i, cond in enumerate(data.conditions, start=ord("a")):
            doc.add_paragraph(f"({chr(i)})  {cond}", style="List Bullet")
    else:
        doc.add_paragraph(PLACEHOLDER)

    # ---- §7 Recommendation ----
    doc.add_heading("7.  Recommendation", level=2)
    doc.add_paragraph(data.recommendation)

    # ---- §8 Human Approval Section ----
    doc.add_heading("8.  Approval Decision", level=2)
    doc.add_paragraph(
        "[ ]  APPROVED for continued operation under the conditions in §6 above.\n"
        "[ ]  REJECTED — immediate corrective action required."
    )
    doc.add_paragraph()
    sig_table = doc.add_table(rows=2, cols=2)
    sig_table.style = "Table Grid"
    sig_table.cell(0, 0).text = f"Area Engineer Signature: {data.area_engineer_name}"
    sig_table.cell(0, 1).text = f"Date: {stamp_date}"
    sig_table.cell(1, 0).text = f"Inspector Signature:      {data.inspector_name}"
    sig_table.cell(1, 1).text = f"Date: {stamp_date}"

    # ---- §9 Evidence Table (citations) ----
    doc.add_heading("9.  Evidence Table", level=2)
    if data.evidence:
        ev_table = doc.add_table(rows=1, cols=5)
        ev_table.style = "Table Grid"
        ev_hdr = ev_table.rows[0].cells
        for cell, txt in zip(
            ev_hdr,
            ["#", "Source File", "Page", "Confidence", "Exact Quote"],
        ):
            cell.text = txt
            _bold_cell(cell)
        for i, ev in enumerate(data.evidence, 1):
            row = ev_table.add_row().cells
            row[0].text = str(i)
            row[1].text = ev.source_file
            row[2].text = str(ev.page_number)
            row[3].text = f"{ev.confidence:.0%}"
            row[4].text = ev.exact_quote[:300]
    else:
        doc.add_paragraph("No RAG evidence was retrieved for this run.")

    # ---- Footer note ----
    doc.add_paragraph()
    footer = doc.add_paragraph(
        f"Generated by Sovereign AI Workbench — {datetime.now(timezone.utc).strftime('%Y-%m-%dT%H:%M:%SZ')} UTC | "
        "Zero external API calls made during generation."
    )
    footer.runs[0].font.size = Pt(8)
    footer.runs[0].font.color.rgb = RGBColor(0x88, 0x88, 0x88)

    out = _output_path(data.run_id, f"approval_note_{note_id}.docx")
    doc.save(str(out))

    logger.info(
        "DOCGEN_DOCX | note_id=%s | equipment=%s | path=%s | "
        "findings=%d | evidence=%d",
        note_id, data.equipment_id, out, len(data.findings), len(data.evidence),
    )
    return out.resolve()


# ---------------------------------------------------------------------------
# XLSX — Calculation / Measurement Sheet
# ---------------------------------------------------------------------------

def generate_calculation_sheet_xlsx(data: CalculationSheetInput) -> Path:
    """Generate an Excel workbook with a measurement-vs-limit table.

    Returns the absolute path to the created file.
    """
    import openpyxl
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Measurements"

    # ---- Palette ----
    HEADER_FILL  = PatternFill("solid", fgColor="1A376E")
    ALERT_FILL   = PatternFill("solid", fgColor="FFD700")
    NORMAL_FILL  = PatternFill("solid", fgColor="C6EFCE")
    CAUTION_FILL = PatternFill("solid", fgColor="FFEB9C")
    TRIP_FILL    = PatternFill("solid", fgColor="FFC7CE")
    thin_side    = Side(style="thin", color="AAAAAA")
    thin_border  = Border(left=thin_side, right=thin_side, top=thin_side, bottom=thin_side)

    _STATUS_FILL = {
        "NORMAL":   NORMAL_FILL,
        "ALERT":    ALERT_FILL,
        "FLAG":     ALERT_FILL,
        "CAUTION":  CAUTION_FILL,
        "SHUTDOWN": TRIP_FILL,
        "REJECT":   TRIP_FILL,
    }

    # ---- Title rows ----
    ws.merge_cells("A1:F1")
    ws["A1"] = data.title
    ws["A1"].font = Font(bold=True, size=14, color="FFFFFF")
    ws["A1"].fill = HEADER_FILL
    ws["A1"].alignment = Alignment(horizontal="center")

    ws.merge_cells("A2:F2")
    ws["A2"] = (
        f"Equipment: {data.equipment_id}    |    "
        f"Generated: {_today()}    |    "
        f"Run ID: {data.run_id or 'N/A'}"
    )
    ws["A2"].font = Font(italic=True, size=10)
    ws["A2"].alignment = Alignment(horizontal="center")

    # ---- Column headers ----
    headers = ["#", "Parameter", "Reading", "Unit", "Limit", "Status"]
    widths   = [5,   36,          16,        10,     26,       12]
    for col, (hdr, width) in enumerate(zip(headers, widths), 1):
        cell = ws.cell(row=3, column=col, value=hdr)
        cell.font      = Font(bold=True, color="FFFFFF")
        cell.fill      = HEADER_FILL
        cell.alignment = Alignment(horizontal="center")
        cell.border    = thin_border
        ws.column_dimensions[get_column_letter(col)].width = width

    # ---- Data rows ----
    for i, m in enumerate(data.measurements, 1):
        row = 3 + i
        status = str(m.get("status", "NORMAL")).upper()
        fill   = _STATUS_FILL.get(status, PatternFill())
        values = [
            i,
            m.get("parameter", PLACEHOLDER),
            m.get("reading",   PLACEHOLDER),
            m.get("unit",      ""),
            m.get("limit",     PLACEHOLDER),
            status,
        ]
        for col, val in enumerate(values, 1):
            cell = ws.cell(row=row, column=col, value=val)
            cell.border = thin_border
            if col > 1:
                cell.fill = fill
            cell.alignment = Alignment(horizontal="center" if col != 2 else "left")

    # ---- Footer ----
    footer_row = 4 + len(data.measurements)
    ws.cell(
        row=footer_row, column=1,
        value=(
            f"Generated by Sovereign AI Workbench — {_today()} | "
            "Zero external API calls."
        ),
    ).font = Font(italic=True, size=8, color="888888")

    ws.freeze_panes = "A4"

    out = _output_path(data.run_id, f"calculations_{data.equipment_id}.xlsx")
    wb.save(str(out))

    logger.info(
        "DOCGEN_XLSX | title=%r | equipment=%s | rows=%d | path=%s",
        data.title, data.equipment_id, len(data.measurements), out,
    )
    return out.resolve()


# ---------------------------------------------------------------------------
# PPTX — Summary Deck
# ---------------------------------------------------------------------------

def generate_summary_deck_pptx(data: SummaryDeckInput) -> Path:
    """Generate a minimal PPTX summary deck.

    Slides: Title → one slide per finding → Recommendation/Conclusion.
    Returns the absolute path to the created file.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = _emu(25.4)   # 25.4 cm wide  (widescreen 16:9 approx)
    prs.slide_height = _emu(14.29)

    DARK_BLUE = RGBColor(0x1A, 0x37, 0x6E)
    AMBER     = RGBColor(0xFF, 0xD7, 0x00)
    WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

    _STATUS_COLOUR = {
        "NORMAL":   RGBColor(0x00, 0x87, 0x54),
        "ALERT":    RGBColor(0xFF, 0xD7, 0x00),
        "FLAG":     RGBColor(0xFF, 0xD7, 0x00),
        "CAUTION":  RGBColor(0xFF, 0x8C, 0x00),
        "SHUTDOWN": RGBColor(0xC0, 0x00, 0x00),
        "REJECT":   RGBColor(0xC0, 0x00, 0x00),
    }

    blank_layout = prs.slide_layouts[6]  # completely blank

    # ---- Slide 1: Title ----
    s = prs.slides.add_slide(blank_layout)
    _bg_colour(s, DARK_BLUE)
    _add_textbox(
        s, data.title,
        l=_emu(1), t=_emu(4), w=_emu(23), h=_emu(3),
        font_size=Pt(32), bold=True, colour=WHITE, align=PP_ALIGN.CENTER,
    )
    _add_textbox(
        s,
        f"Equipment: {data.equipment_id}   |   Report: {data.inspection_report_id}\n"
        f"Sovereign AI Workbench   |   {_today()}",
        l=_emu(1), t=_emu(7.5), w=_emu(23), h=_emu(2),
        font_size=Pt(16), colour=AMBER, align=PP_ALIGN.CENTER,
    )

    # ---- Slide(s): one per finding ----
    for f in data.findings:
        s = prs.slides.add_slide(blank_layout)
        status_colour = _STATUS_COLOUR.get(f.status.upper(), WHITE)
        _bg_colour(s, RGBColor(0xF5, 0xF7, 0xFA))

        # Coloured header strip
        _add_coloured_rect(s, colour=DARK_BLUE, l=0, t=0, w=_emu(25.4), h=_emu(2))
        _add_textbox(
            s, f"{f.label}",
            l=_emu(0.5), t=_emu(0.2), w=_emu(18), h=_emu(1.5),
            font_size=Pt(22), bold=True, colour=WHITE,
        )
        # Status badge (top-right)
        _add_textbox(
            s, f.status,
            l=_emu(18.5), t=_emu(0.3), w=_emu(6), h=_emu(1.3),
            font_size=Pt(18), bold=True, colour=status_colour, align=PP_ALIGN.CENTER,
        )

        # Finding description
        _add_textbox(
            s, f.description,
            l=_emu(0.8), t=_emu(2.4), w=_emu(23.8), h=_emu(5),
            font_size=Pt(16), colour=RGBColor(0x1A, 0x1A, 0x1A),
        )

        # Measurement row
        meas_parts = []
        if f.measured_value:
            meas_parts.append(f"Measured:  {f.measured_value}")
        if f.limit_value:
            meas_parts.append(f"Limit:  {f.limit_value}")
        if meas_parts:
            _add_textbox(
                s, "   |   ".join(meas_parts),
                l=_emu(0.8), t=_emu(7.8), w=_emu(23.8), h=_emu(1.5),
                font_size=Pt(14), colour=status_colour, bold=True,
            )

    # ---- Last slide: Recommendation ----
    s = prs.slides.add_slide(blank_layout)
    _bg_colour(s, DARK_BLUE)
    _add_textbox(
        s, "Recommendation",
        l=_emu(1), t=_emu(0.8), w=_emu(23), h=_emu(1.5),
        font_size=Pt(26), bold=True, colour=AMBER,
    )
    _add_textbox(
        s, data.recommendation,
        l=_emu(1), t=_emu(2.5), w=_emu(23), h=_emu(8),
        font_size=Pt(16), colour=WHITE,
    )
    _add_textbox(
        s,
        f"Generated by Sovereign AI Workbench | {_today()} | Zero external API calls.",
        l=_emu(1), t=_emu(12.5), w=_emu(23), h=_emu(0.8),
        font_size=Pt(9), colour=RGBColor(0x99, 0xAA, 0xBB),
    )

    out = _output_path(data.run_id, f"summary_{data.equipment_id}.pptx")
    prs.save(str(out))

    logger.info(
        "DOCGEN_PPTX | title=%r | equipment=%s | slides=%d | path=%s",
        data.title, data.equipment_id, len(prs.slides), out,
    )
    return out.resolve()


# ---------------------------------------------------------------------------
# Convenience factory — called by the planner with a flat dict
# ---------------------------------------------------------------------------

def generate_from_agent_output(
    doc_format: str,
    run_id: Optional[int],
    payload: dict,
) -> Path:
    """Dispatch to the correct generator based on *doc_format*.

    Args:
        doc_format: ``"docx"`` | ``"xlsx"`` | ``"pptx"``
        run_id:     AgentRun.id (used for the output sub-folder)
        payload:    Dict matching the corresponding ``*Input`` dataclass fields.

    Returns the absolute path to the generated file.

    Raises ``ValueError`` for unknown formats, ``TypeError`` if required
    fields are missing from *payload*.
    """
    fmt = doc_format.lower().strip(".")

    if fmt == "docx":
        return generate_approval_note_docx(ApprovalNoteInput(**payload))

    if fmt == "xlsx":
        return generate_calculation_sheet_xlsx(CalculationSheetInput(**payload))

    if fmt == "pptx":
        return generate_summary_deck_pptx(SummaryDeckInput(**payload))

    raise ValueError(
        f"Unknown doc_format {doc_format!r}. "
        "Supported formats: 'docx', 'xlsx', 'pptx'."
    )


# ---------------------------------------------------------------------------
# Internal DOCX helpers
# ---------------------------------------------------------------------------

def _cm(cm: float):
    from docx.shared import Cm
    return Cm(cm)


def _add_label_value(doc, label: str, value: str) -> None:
    """Add a single-paragraph ``Label:  Value`` line with the label in bold."""
    from docx.shared import Pt
    p = doc.add_paragraph()
    run_label = p.add_run(f"{label}:  ")
    run_label.bold = True
    run_label.font.size = Pt(10)
    run_value = p.add_run(value)
    run_value.font.size = Pt(10)
    p.paragraph_format.space_after = Pt(0)  # space_after is None by default; use Pt directly


def _bold_cell(cell) -> None:
    """Make all runs in a table cell bold."""
    for para in cell.paragraphs:
        for run in para.runs:
            run.bold = True


# ---------------------------------------------------------------------------
# Internal PPTX helpers
# ---------------------------------------------------------------------------

def _emu(cm_val: float) -> int:
    """Convert centimetres to EMU (English Metric Units) for python-pptx."""
    from pptx.util import Cm
    return Cm(cm_val)


def _bg_colour(slide, colour) -> None:
    """Fill the slide background with a solid colour."""
    from pptx.util import Pt
    from pptx.oxml.ns import qn
    from lxml import etree

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _add_textbox(slide, text: str, *, l, t, w, h, font_size, colour, bold=False, align=None):
    """Add a textbox to *slide* and style it."""
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.color.rgb = colour
    run.font.bold = bold
    return txBox


def _add_coloured_rect(slide, *, colour, l, t, w, h):
    """Add a filled rectangle shape (used as a coloured banner)."""
    from pptx.util import Pt
    shape = slide.shapes.add_shape(
        1,   # MSO_SHAPE_TYPE.RECTANGLE
        l, t, w, h,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = colour
    shape.line.fill.background()  # no border
    return shape
